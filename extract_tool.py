import os
import sys
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
import zipfile

class ImageExtractorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Office 图片提取工具")
        self.root.geometry("500x350")
        
        # 变量
        self.file_path = tk.StringVar()
        self.output_dir = tk.StringVar()
        
        # UI 布局
        self.create_widgets()
        
    def create_widgets(self):
        # 顶部说明
        header_label = tk.Label(self.root, text="支持 PPT (.pptx) 和 Word (.docx) 文件", font=("Arial", 10))
        header_label.pack(pady=10)
        
        # 文件选择区域
        input_frame = tk.LabelFrame(self.root, text="第一步: 选择文件", padx=10, pady=10)
        input_frame.pack(fill="x", padx=10, pady=5)
        
        tk.Entry(input_frame, textvariable=self.file_path, width=40).pack(side=tk.LEFT, padx=5)
        tk.Button(input_frame, text="浏览...", command=self.select_file).pack(side=tk.LEFT)
        
        # 输出目录选择区域
        output_frame = tk.LabelFrame(self.root, text="第二步: 选择保存位置", padx=10, pady=10)
        output_frame.pack(fill="x", padx=10, pady=5)
        
        tk.Entry(output_frame, textvariable=self.output_dir, width=40).pack(side=tk.LEFT, padx=5)
        tk.Button(output_frame, text="浏览...", command=self.select_output_dir).pack(side=tk.LEFT)
        
        # 提取按钮
        self.extract_btn = tk.Button(self.root, text="开始提取图片", command=self.start_extraction, 
                                   bg="#0078D7", fg="white", font=("Arial", 12, "bold"), height=2)
        self.extract_btn.pack(fill="x", padx=20, pady=20)
        
        # 状态栏
        self.status_label = tk.Label(self.root, text="准备就绪", bd=1, relief=tk.SUNKEN, anchor=tk.W)
        self.status_label.pack(side=tk.BOTTOM, fill=tk.X)

    def select_file(self):
        file_selected = filedialog.askopenfilename(
            filetypes=[("Office Files", "*.pptx;*.docx"), ("PowerPoint", "*.pptx"), ("Word", "*.docx")]
        )
        if file_selected:
            self.file_path.set(file_selected)
            # 默认输出目录设为文件所在目录
            if not self.output_dir.get():
                self.output_dir.set(os.path.dirname(file_selected))

    def select_output_dir(self):
        dir_selected = filedialog.askdirectory()
        if dir_selected:
            self.output_dir.set(dir_selected)

    def log(self, message):
        self.status_label.config(text=message)
        self.root.update()

import threading

    def start_extraction(self):
        source_file = self.file_path.get()
        target_dir = self.output_dir.get()
        
        if not source_file or not os.path.exists(source_file):
            messagebox.showerror("错误", "请选择有效的源文件！")
            return
            
        if not target_dir:
            messagebox.showerror("错误", "请选择保存位置！")
            return
            
        # 禁用按钮防止重复点击
        self.extract_btn.config(state=tk.DISABLED, text="正在处理...")
        
        # 在新线程中运行提取任务，避免界面卡顿
        threading.Thread(target=self.run_extraction_task, args=(source_file, target_dir), daemon=True).start()

    def run_extraction_task(self, source_file, target_dir):
        try:
            # 创建存放图片的子文件夹
            filename = os.path.splitext(os.path.basename(source_file))[0]
            images_dir = os.path.join(target_dir, f"{filename}_images")
            
            if not os.path.exists(images_dir):
                os.makedirs(images_dir)
            
            count = 0
            if source_file.lower().endswith(".pptx"):
                count = self.extract_from_ppt(source_file, images_dir)
            elif source_file.lower().endswith(".docx"):
                count = self.extract_from_word(source_file, images_dir)
            else:
                self.root.after(0, messagebox.showerror, "错误", "不支持的文件格式！目前仅支持 .pptx 和 .docx")
                return

            if count > 0:
                self.root.after(0, lambda: messagebox.showinfo("成功", f"提取完成！\n共保存 {count} 张图片到：\n{images_dir}"))
                self.root.after(0, lambda: os.startfile(images_dir))
            else:
                self.root.after(0, messagebox.showinfo, "提示", "未在文件中找到图片。")
                
        except Exception as e:
            self.root.after(0, messagebox.showerror, "错误", f"发生错误：{str(e)}")
        finally:
            self.root.after(0, self.reset_ui)

    def reset_ui(self):
        self.extract_btn.config(state=tk.NORMAL, text="开始提取图片")
        self.log("准备就绪")

    def extract_from_ppt(self, ppt_path, output_folder):
        prs = Presentation(ppt_path)
        image_count = 0
        
        for i, slide in enumerate(prs.slides):
            self.log(f"正在处理幻灯片 {i+1}/{len(prs.slides)}...")
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_ext = image.ext
                    image_filename = f"Slide{i+1}_Img{image_count+1}.{image_ext}"
                    image_path = os.path.join(output_folder, image_filename)
                    
                    with open(image_path, 'wb') as f:
                        f.write(image.blob)
                    image_count += 1
        return image_count

    def extract_from_word(self, docx_path, output_folder):
        """
        使用 zipfile 方式提取 Word 中的所有图片，这种方式最全面，
        虽然顺序可能不是完全严格的阅读顺序，但通常是按照文件内部ID排序的。
        """
        image_count = 0
        self.log("正在分析 Word 文档结构...")
        
        with zipfile.ZipFile(docx_path) as z:
            all_files = z.namelist()
            # 过滤出媒体文件夹中的图片
            media_files = [f for f in all_files if f.startswith("word/media/")]
            
            for media_file in media_files:
                # 获取扩展名
                ext = os.path.splitext(media_file)[1]
                if ext.lower() not in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
                    continue
                    
                image_count += 1
                self.log(f"正在提取第 {image_count} 张图片...")
                
                # 读取图片数据
                image_data = z.read(media_file)
                
                # 保存图片
                # 使用原始文件名中的数字部分来保持一定的顺序，或者直接重命名
                original_name = os.path.basename(media_file)
                target_path = os.path.join(output_folder, f"Image_{image_count}{ext}")
                
                with open(target_path, 'wb') as f:
                    f.write(image_data)
                    
        return image_count

if __name__ == "__main__":
    root = tk.Tk()
    app = ImageExtractorApp(root)
    root.mainloop()
